import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()
    
    # 1. Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "SYNC-HUB 시스템 기획 보고서"
    subtitle.text = "데이터 중앙수집 및 업무 자동화를 통한 효율성 증대\n\n작성일자: 2026-02-21"

    # Slide 2: AS-IS 1
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    slide.shapes.title.text = "AS-IS 1. 프로젝트 검토 단계의 진행"
    tf = slide.placeholders[1].text_frame
    tf.text = "영업팀&고객의뢰로 검토 진행"
    tf.add_paragraph().text = "테스트 진행 및 컨셉 설계 진행"
    tf.add_paragraph().text = "테스트 결과와 컨셉 설계 기준으로 제안서 및 예산 내역 작성"
    tf.add_paragraph().text = "영업팀에서 고객사에 견적 제출"
    p = tf.add_paragraph()
    p.text = "산출물: 제안서, 원가검토, 기구컨셉도면, 유틸사양, 회의록, 일정표, 견적서 등"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "고객사 자료: 1차 사양서, 의뢰서, 설치 장소 레이아웃 등"
    p.level = 1

    # Slide 3: AS-IS 2
    slide = prs.slides.add_slide(bullet_slide_layout)
    slide.shapes.title.text = "AS-IS 2. 프로젝트 제작 단계의 진행"
    tf = slide.placeholders[1].text_frame
    tf.text = "프로젝트 수주 완료 후 제작 시작"
    tf.add_paragraph().text = "기구/전장/SW 설계 진행"
    tf.add_paragraph().text = "설계 완료 후 파츠 발주 진행"
    tf.add_paragraph().text = "파츠 입고 후 기구/전장 조립"
    tf.add_paragraph().text = "조립 완료 후 테스트 진행"
    tf.add_paragraph().text = "출하 진행"
    p = tf.add_paragraph()
    p.text = "산출물: 파츠 입고 관리 대장, 집행 금액 내역, 기구/전장 도면, SW 소스코드, 펀치리스트 등"
    p.level = 1

    # Slide 4: AS-IS 3
    slide = prs.slides.add_slide(bullet_slide_layout)
    slide.shapes.title.text = "AS-IS 3. 프로젝트 설치 단계의 진행"
    tf = slide.placeholders[1].text_frame
    tf.text = "반입 후 기구/전장 설치"
    tf.add_paragraph().text = "설비 테스트 진행"
    tf.add_paragraph().text = "고객사 설비 교육 진행"
    p = tf.add_paragraph()
    p.text = "산출물: 펀치리스트, 교육 진행 관리 대장, 집행 금액 내역 등"
    p.level = 1

    # Slide 5: 기획 의도
    slide = prs.slides.add_slide(bullet_slide_layout)
    slide.shapes.title.text = "기획 의도 및 기대 효과"
    tf = slide.placeholders[1].text_frame
    tf.text = "파편화 된 자료와 데이터들을 중앙수집"
    tf.add_paragraph().text = "수집 된 데이터 검색 시스템으로 업무 속도 상승 (글로벌 검색 및 리스트 검색)"
    tf.add_paragraph().text = "정확한 데이터 기반으로 프로젝트 진행 (사양서와 안건)"
    tf.add_paragraph().text = "수집 된 데이터를 기반으로 자동화 시스템 구현 (EXCEL, PDF 내보내기)"
    tf.add_paragraph().text = "일정한 양식 기반 산출물로 인원간 원활한 인수인계"
    tf.add_paragraph().text = "미래: AI 활용 데이터 생성/분석으로 프로젝트 성공율 제고"

    def add_image_slide(title_text, desc_lines, img_paths):
        layout = prs.slide_layouts[5] # Title only
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = title_text
        
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9), Inches(1.5))
        tf = txBox.text_frame
        for i, line in enumerate(desc_lines):
            if i == 0:
                tf.text = line
            else:
                tf.add_paragraph().text = line

        num_imgs = len(img_paths)
        if num_imgs == 1:
            try:
                slide.shapes.add_picture(img_paths[0], Inches(1.5), Inches(2.7), height=Inches(4.2))
            except Exception as e:
                print(f"Error loading {img_paths[0]}: {e}")
        elif num_imgs == 2:
            try:
                slide.shapes.add_picture(img_paths[0], Inches(0.5), Inches(2.7), height=Inches(3.8))
                slide.shapes.add_picture(img_paths[1], Inches(5.0), Inches(2.7), height=Inches(3.8))
            except:
                pass
        elif num_imgs == 3:
            try:
                slide.shapes.add_picture(img_paths[0], Inches(0.2), Inches(3.0), height=Inches(3.0))
                slide.shapes.add_picture(img_paths[1], Inches(3.5), Inches(3.0), height=Inches(3.0))
                slide.shapes.add_picture(img_paths[2], Inches(6.8), Inches(3.0), height=Inches(3.0))
            except:
                pass
        else:
            for i, ip in enumerate(img_paths[:4]):
                try:
                    left = Inches(0.5 + (i%2)*4.8)
                    top = Inches(2.8 + (i//2)*2.3)
                    slide.shapes.add_picture(ip, left, top, height=Inches(2.0))
                except:
                    pass

    base_dir = "/home/arari123/웹 스크린샷"
    
    # Slide 6
    add_image_slide("기능 소개 1. 글로벌 검색 기능", 
                    ["- 모든 자료들을 자연어로 검색할 수 있는 글로벌 검색 기능"], 
                    [os.path.join(base_dir, "글로벌 검색_메인.png"),
                     os.path.join(base_dir, "글로벌 검색_자료 검색 결과.png")])

    # Slide 7
    add_image_slide("기능 소개 2. 프로젝트 리스트 페이지",
                    ["- 모든 프로젝트를 가시성있게 표현하여 프로젝트 현황 파악", 
                     "- 변경/추가사항 발생 시 알람 표시로 빠른 추적"],
                    [os.path.join(base_dir, "메인 페이지_설비 프로젝트 리스트.png"),
                     os.path.join(base_dir, "메인 페이지_파츠 프로젝트 리스트.png")])

    # Slide 8
    add_image_slide("기능 소개 3. 프로젝트 메인 페이지",
                    ["- 프로젝트의 전체상황을 한 페이지내에서 요약 정리"],
                    [os.path.join(base_dir, "프로젝트 메인_메인.png")])

    # Slide 9
    add_image_slide("기능 소개 4. 예산 관리",
                    ["- 친숙한 디자인, 전체 예산/집행 추적, 재료/인건/경비 상세 내역 관리",
                     "- 향후 결재 보고서 내보내기, ERP 금액 업로드, 견적서 자동작성 추가"],
                    [os.path.join(base_dir, "예산 페이지_메인.png"),
                     os.path.join(base_dir, "예산 페이지_경비.png")])

    # Slide 10
    add_image_slide("기능 소개 5. 일정 관리",
                    ["- 프로젝트에 대한 일정 조회 및 작성 (기본 양식으로 내보내기 지원)"],
                    [os.path.join(base_dir, "일정 페이지_일정 관리.png") if not os.path.exists(os.path.join(base_dir, "일정 페이지_일정 조회.png")) else os.path.join(base_dir, "일정 페이지_일정 조회.png"),
                     os.path.join(base_dir, "일정 페이지_일정 작성.png")])

    # Slide 11
    add_image_slide("기능 소개 6. 안건 관리",
                    ["- ISSUE, Q&A, TODO, 회의록 등 아웃룩 UI 모방 디자인"],
                    [os.path.join(base_dir, "안건 페이지_조회.png"),
                     os.path.join(base_dir, "안건 페이지_작업보고서 작성.png")])

    # Slide 12
    add_image_slide("기능 소개 7. 데이터 관리",
                    ["- 자료 업로드 및 파일 다운로드, 글로벌 검색 색인 연동"],
                    [os.path.join(base_dir, "데이터 관리 페이지.png")])

    # Slide 13
    add_image_slide("기능 소개 8. (임시) AI 답변 생성",
                    ["- GEMINI 2.5 FLASH LITE 등 모델 연동 (문서 요약, 안건 요약)"],
                    [os.path.join(base_dir, "RAG 임시 구현 페이지.png")])

    # Slide 14
    slide = prs.slides.add_slide(bullet_slide_layout)
    slide.shapes.title.text = "기능 소개 9. (구상 중) 설비 사양서 작성"
    tf = slide.placeholders[1].text_frame
    tf.text = "프로젝트 검토 단계부터 사양서를 작성하여 누락되는 예산 방지"
    tf.add_paragraph().text = "고객/협업자와 변경사항 기록 공유로 분쟁 최소화"
    tf.add_paragraph().text = "스마트 위저드 UI를 사용하여 꼼꼼한 체크 및 입력 지원"
    tf.add_paragraph().text = "기존 데이터 활용하여 불필요한 입력 방지"
    tf.add_paragraph().text = "향후 AI 고도화를 통해 더 쉬운 작성 지원"

    out_path = "/home/arari123/SYNC-HUB_기획보고서.pptx"
    prs.save(out_path)
    print(f"Presentation saved successfully to {out_path}")

if __name__ == '__main__':
    create_presentation()
