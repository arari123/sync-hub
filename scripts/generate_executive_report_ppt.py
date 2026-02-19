#!/usr/bin/env python3
"""Generate executive planning PPT from context + screenshots.

Run in Docker:
  docker exec -w /app synchub_web python3 scripts/generate_executive_report_ppt.py
"""

from __future__ import annotations

import argparse
import json
from datetime import datetime
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


NAVY = RGBColor(11, 37, 69)
CYAN = RGBColor(20, 122, 169)
DARK = RGBColor(31, 41, 55)
GRAY = RGBColor(75, 85, 99)
WHITE = RGBColor(255, 255, 255)
GREEN = RGBColor(5, 150, 105)
ORANGE = RGBColor(217, 119, 6)


def _read_json(path: Path) -> dict:
    return json.loads(path.read_text(encoding="utf-8"))


def _fmt_money(value: int) -> str:
    return f"{int(value):,}원"


def _add_title(slide, title: str, subtitle: str = "") -> None:
    bar = slide.shapes.add_shape(
        autoshape_type_id=1,  # rectangle
        left=Inches(0),
        top=Inches(0),
        width=Inches(13.333),
        height=Inches(0.58),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.55), Inches(0.7), Inches(12.2), Inches(0.7))
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.bold = True
    p.font.size = Pt(30)
    p.font.color.rgb = NAVY

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.58), Inches(1.25), Inches(12.2), Inches(0.38))
        tf2 = sub_box.text_frame
        tf2.clear()
        p2 = tf2.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(14)
        p2.font.color.rgb = GRAY


def _add_bullets(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    bullets: list[str],
    font_size: int = 16,
    color: RGBColor = DARK,
) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.clear()
    for index, line in enumerate(bullets):
        p = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        p.text = f"• {line}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.space_after = Pt(7)


def _add_text(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    text: str,
    font_size: int = 13,
    color: RGBColor = DARK,
    bold: bool = False,
    align=PP_ALIGN.LEFT,
) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bool(bold)


def _capture_path(captures: dict, key: str) -> Path | None:
    item = captures.get(key)
    if not item:
        return None
    path = Path(str(item.get("file", ""))).resolve()
    return path if path.exists() else None


def _add_image(
    slide,
    captures: dict,
    key: str,
    left: float,
    top: float,
    width: float,
    caption: str = "",
) -> None:
    path = _capture_path(captures, key)
    if path:
        slide.shapes.add_picture(str(path), Inches(left), Inches(top), width=Inches(width))
    else:
        _add_text(
            slide,
            left,
            top + 0.8,
            width,
            0.5,
            f"[이미지 누락] {key}",
            font_size=12,
            color=ORANGE,
            bold=True,
            align=PP_ALIGN.CENTER,
        )

    if caption:
        _add_text(
            slide,
            left,
            top + 3.55,
            width,
            0.3,
            caption,
            font_size=11,
            color=GRAY,
            align=PP_ALIGN.CENTER,
        )


def _chunk(items: list, size: int) -> list[list]:
    return [items[index:index + size] for index in range(0, len(items), size)]


def build_ppt(context: dict, manifest: dict, output_path: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    captures = {item.get("key"): item for item in manifest.get("captures", [])}
    stats = context.get("stats", {})
    scenarios = context.get("scenarios", {})
    reps = context.get("representative_projects", {})
    page_inventory = context.get("page_inventory", [])

    generated_date = datetime.now().strftime("%Y-%m-%d")

    # 1) Cover
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(
        slide,
        "Sync-Hub 경영진 보고용 기획서",
        "프로젝트 전체 사항에 대한 효율적 관리 웹앱 현황 및 도입 효과",
    )
    _add_text(slide, 0.58, 2.0, 12.2, 0.5, f"작성일: {generated_date}", font_size=14, color=GRAY)
    _add_text(
        slide,
        0.58,
        2.45,
        12.2,
        0.7,
        f"현재 데이터 기준: 프로젝트 {stats.get('projects_total', 0)}건 / 안건 {stats.get('agenda_threads_total', 0)}건 / 자료 {stats.get('documents_total', 0)}건",
        font_size=18,
        color=CYAN,
        bold=True,
    )
    _add_bullets(
        slide,
        0.78,
        3.25,
        12.0,
        2.9,
        [
            "파편화된 프로젝트 데이터(예산/안건/일정/자료)를 단일 웹 플랫폼으로 통합",
            "검색 기반 접근성 향상으로 개인/팀의 업무 처리 속도 개선",
            "이력 데이터 기반 자동화(EXCEL/PDF/PPT) 및 AI 확장 기반 마련",
        ],
        font_size=18,
    )

    # 2) Executive Summary
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(slide, "Executive Summary", "현재 구현 범위와 경영진 관점 핵심 포인트")
    _add_bullets(
        slide,
        0.65,
        1.85,
        12.0,
        4.2,
        [
            "현재 웹앱은 프로젝트 생성부터 예산/안건/일정/자료관리까지 운영 가능한 수준으로 구현됨",
            "프로젝트 유형(설비/파츠/AS)별 시나리오를 동일한 데이터 모델로 통합 관리함",
            "프로젝트 자료실 + 글로벌 검색 + 데이터허브를 통해 문서 검색성과 활용성을 높임",
            "현재 누적된 데이터는 향후 자동 문서 생성 및 AI 분석 고도화의 기초 자산이 됨",
        ],
        font_size=18,
    )
    _add_text(
        slide,
        0.7,
        6.35,
        12.0,
        0.5,
        "주요 지표: "
        f"프로젝트 {stats.get('projects_total', 0)} / 예산버전 {stats.get('budget_versions_total', 0)} / "
        f"안건엔트리 {stats.get('agenda_entries_total', 0)} / 자료폴더 {stats.get('document_folders_total', 0)}",
        font_size=12,
        color=GRAY,
    )

    # 3) 기획의도
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(slide, "기획의도", "프로젝트 전체 사항에 대한 효율적 관리를 위한 웹앱 개발")
    _add_bullets(
        slide,
        0.75,
        1.8,
        12.0,
        4.8,
        [
            "파편화된 자료/데이터를 중앙 수집하여 DB화하고 향후 자동화에 활용",
            "수집 데이터 검색 시스템으로 개인의 일처리 속도 상승",
            "사양서/안건 기반의 정확한 데이터 중심 프로젝트 진행",
            "수집 데이터를 활용한 자동화(EXCEL, PDF, PPT 내보내기)로 업무 효율 증대",
            "향후 AI 기반 데이터 생성/분석/보완으로 프로젝트 수행 속도와 성공률 극대화",
        ],
        font_size=17,
    )

    # 4) AS-IS vs TO-BE (1)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(slide, "AS-IS vs TO-BE (1)", "검토 단계와 Drop 프로젝트 데이터 활용")
    _add_text(slide, 0.75, 1.75, 6.0, 0.36, "AS-IS", font_size=16, bold=True, color=ORANGE)
    _add_bullets(
        slide,
        0.75,
        2.05,
        6.0,
        3.9,
        [
            "PPT/Excel 중심 수작업 제안서·원가검토",
            "부서별 상이한 양식으로 메일 협업 부담",
            "검토 단계 Drop 시 자료는 남지만 사실상 재활용 어려움",
        ],
        font_size=14,
    )
    _add_text(slide, 6.95, 1.75, 5.6, 0.36, "TO-BE", font_size=16, bold=True, color=GREEN)
    _add_bullets(
        slide,
        6.95,
        2.05,
        5.6,
        3.9,
        [
            "검토 단계부터 웹에 데이터 축적(사양/원가/이력)",
            "검토 히스토리 재활용으로 유사 프로젝트 대응 속도 향상",
            "중앙 저장소 기반 조직 차원의 지식 자산화",
        ],
        font_size=14,
    )

    # 5) AS-IS vs TO-BE (2)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(slide, "AS-IS vs TO-BE (2)", "설계~설치 협업, 자료검색, 문서 양식 통일")
    _add_bullets(
        slide,
        0.75,
        1.85,
        12.0,
        4.8,
        [
            "AS-IS: 설계~설치 이력이 개인별 Excel/PPT로 분산되어 취합 비용이 큼",
            "TO-BE: 웹에서 단일 이력 관리로 협업 취합 작업을 구조적으로 제거",
            "AS-IS: 로컬/메일/부서 폴더 중심으로 자료 탐색 난이도 높음",
            "TO-BE: 웹 검색(프로젝트/안건/문서)으로 자료 접근 속도 향상",
            "AS-IS: PM별 상이한 원가검토 양식으로 문서 해석 난이도 발생",
            "TO-BE: 통일 양식 기반으로 누구나 조회/해석 가능한 표준 체계 확립",
        ],
        font_size=15,
    )

    # 6) 유형별 데이터 현황
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(slide, "현재 데이터 현황", "유형별 프로젝트 분포(시드 데이터 기준)")
    type_counts = stats.get("projects_by_type", {})
    stage_counts = stats.get("projects_by_stage", {})
    _add_text(
        slide,
        0.8,
        1.85,
        5.8,
        0.5,
        f"프로젝트 총 {stats.get('projects_total', 0)}건",
        font_size=24,
        bold=True,
        color=NAVY,
    )
    _add_bullets(
        slide,
        0.95,
        2.4,
        5.5,
        2.6,
        [
            f"설비: {type_counts.get('equipment', 0)}건",
            f"파츠: {type_counts.get('parts', 0)}건",
            f"AS: {type_counts.get('as', 0)}건",
            f"안건: {stats.get('agenda_threads_total', 0)}건 / 코멘트: {stats.get('agenda_comments_total', 0)}건",
            f"자료 문서: {stats.get('documents_total', 0)}건 / 폴더: {stats.get('document_folders_total', 0)}건",
        ],
        font_size=15,
    )
    _add_text(slide, 6.65, 1.9, 5.9, 0.4, "진행 단계 분포", font_size=16, bold=True, color=CYAN)
    stage_lines = [f"{key}: {value}건" for key, value in sorted(stage_counts.items())]
    _add_bullets(slide, 6.75, 2.35, 5.7, 3.5, stage_lines, font_size=14)

    # 7) 유형별 시나리오 - 설비
    eq = reps.get("equipment", {}) or {}
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(slide, "유형별 시나리오: 설비 프로젝트", "대표 프로젝트와 운영 흐름")
    _add_bullets(
        slide,
        0.7,
        1.75,
        5.8,
        2.2,
        [
            scenarios.get("equipment", {}).get("summary", ""),
            f"대표: {eq.get('code', '')} {eq.get('name', '')}",
            f"단계: {eq.get('current_stage_label', eq.get('current_stage', ''))}",
            f"예산총액: {_fmt_money((eq.get('metrics') or {}).get('budget_total_krw', 0))}",
        ],
        font_size=14,
    )
    _add_image(slide, captures, "equipment_overview", 0.72, 3.15, 11.7, "프로젝트 메인(설비형)")

    # 8) 유형별 시나리오 - 파츠
    parts = reps.get("parts", {}) or {}
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(slide, "유형별 시나리오: 파츠 프로젝트", "대표 프로젝트와 운영 흐름")
    _add_bullets(
        slide,
        0.7,
        1.75,
        5.8,
        2.2,
        [
            scenarios.get("parts", {}).get("summary", ""),
            f"대표: {parts.get('code', '')} {parts.get('name', '')}",
            f"단계: {parts.get('current_stage_label', parts.get('current_stage', ''))}",
            f"예산총액: {_fmt_money((parts.get('metrics') or {}).get('budget_total_krw', 0))}",
        ],
        font_size=14,
    )
    _add_image(slide, captures, "parts_overview", 0.72, 3.15, 11.7, "프로젝트 메인(파츠형)")

    # 9) 유형별 시나리오 - AS
    as_proj = reps.get("as", {}) or {}
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(slide, "유형별 시나리오: AS 프로젝트", "대표 프로젝트와 운영 흐름")
    _add_bullets(
        slide,
        0.7,
        1.75,
        5.8,
        2.2,
        [
            scenarios.get("as", {}).get("summary", ""),
            f"대표: {as_proj.get('code', '')} {as_proj.get('name', '')}",
            f"단계: {as_proj.get('current_stage_label', as_proj.get('current_stage', ''))}",
            f"연결 설비 프로젝트 ID: {as_proj.get('parent_project_id')}",
        ],
        font_size=14,
    )
    _add_image(slide, captures, "as_overview", 0.72, 3.12, 5.7, "프로젝트 메인(AS형)")
    _add_image(slide, captures, "as_schedule_management", 6.6, 3.12, 5.7, "일정 관리(AS형 안내)")

    # 10+) 페이지 인벤토리 (모든 웹페이지 기능설명)
    inventory_chunks = _chunk(page_inventory, 7)
    total_chunks = len(inventory_chunks)
    for idx, chunk in enumerate(inventory_chunks, start=1):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _add_title(
            slide,
            f"웹페이지 기능 인벤토리 ({idx}/{total_chunks})",
            "현재 구현된 페이지의 기능 설명",
        )
        lines = []
        for row in chunk:
            features = ", ".join(row.get("features", []))
            status = "구현" if row.get("status") == "implemented" else "예정/플레이스홀더"
            lines.append(f"[{status}] {row.get('name', '')} ({row.get('route', '')})")
            lines.append(f"  - {features}")
        _add_bullets(slide, 0.72, 1.75, 12.0, 4.9, lines, font_size=11)

    # Feature slides with screenshots
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(slide, "인증/접근 관리", "로그인, 회원가입, 이메일 인증")
    _add_image(slide, captures, "login", 0.72, 1.95, 4.1, "로그인")
    _add_image(slide, captures, "signup", 4.95, 1.95, 4.1, "회원가입")
    _add_image(slide, captures, "verify_email", 9.18, 1.95, 4.1, "이메일 인증")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(slide, "통합 검색/홈", "프로젝트, 안건, 문서 검색과 요약 대시보드")
    _add_image(slide, captures, "home", 0.72, 1.95, 11.8, "홈/통합검색")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(slide, "프로젝트 생성/설정", "생성 후 운영 정보 관리")
    _add_image(slide, captures, "project_create", 0.72, 1.95, 5.75, "프로젝트 생성")
    _add_image(slide, captures, "project_info_edit", 6.58, 1.95, 5.75, "프로젝트 설정")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(slide, "예산 관리", "예산 메인 + 입력 화면")
    _add_image(slide, captures, "budget_main", 0.72, 1.78, 11.8, "예산 메인")
    _add_image(slide, captures, "budget_edit_material", 0.72, 5.35, 3.7, "재료비 입력")
    _add_image(slide, captures, "budget_edit_labor", 4.84, 5.35, 3.7, "인건비 입력")
    _add_image(slide, captures, "budget_edit_expense", 8.95, 5.35, 3.7, "경비 입력")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(slide, "안건 관리", "목록/작성/상세(답변·코멘트)")
    _add_image(slide, captures, "agenda_list", 0.72, 1.95, 4.1, "안건 목록")
    _add_image(slide, captures, "agenda_create", 4.95, 1.95, 4.1, "안건 작성")
    _add_image(slide, captures, "agenda_detail", 9.18, 1.95, 4.1, "안건 상세")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(slide, "일정 관리", "일정 조회/편집 분리 운영")
    _add_image(slide, captures, "schedule_management", 0.72, 1.95, 5.75, "일정 관리")
    _add_image(slide, captures, "schedule_write", 6.58, 1.95, 5.75, "일정 작성")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(slide, "자료/데이터 관리", "프로젝트 자료실 + 데이터 허브")
    _add_image(slide, captures, "project_data", 0.72, 1.95, 5.75, "데이터 관리(프로젝트 자료실)")
    _add_image(slide, captures, "data_hub", 6.58, 1.95, 5.75, "데이터 허브")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(slide, "사양 관리 및 향후 확장", "현재 상태와 다음 단계")
    _add_image(slide, captures, "spec_placeholder", 0.72, 1.95, 5.7, "사양 관리(현재 플레이스홀더)")
    _add_bullets(
        slide,
        6.55,
        2.15,
        6.0,
        3.6,
        [
            "사양 관리 화면은 현재 안내 페이지로 제공 중",
            "향후 사양 입력-변경이력-자동 문서생성(PPT/PDF) 연계 예정",
            "프로젝트 데이터 누적 기반 AI 분석/추천 기능으로 확장",
        ],
        font_size=15,
    )

    # Closing
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(slide, "기대 효과", "운영 정착 시 기대되는 경영 성과")
    _add_bullets(
        slide,
        0.8,
        1.9,
        12.0,
        4.5,
        [
            "검색/조회 시간 단축을 통한 PM·협업 부서 생산성 향상",
            "프로젝트 수명주기 이력의 표준화로 품질과 재사용성 개선",
            "문서 자동화와 AI 고도화를 위한 데이터 자산 확보",
            "데이터 기반 의사결정으로 납기/원가/리스크 관리 정밀도 향상",
        ],
        font_size=19,
    )
    _add_text(
        slide,
        0.8,
        6.35,
        12.0,
        0.5,
        "Sync-Hub는 프로젝트 운영 체계를 데이터 중심으로 전환하는 실행 플랫폼입니다.",
        font_size=14,
        color=CYAN,
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument(
        "--context",
        default="reports/executive/2026-02-19/context.json",
        help="Context json path",
    )
    parser.add_argument(
        "--manifest",
        default="reports/executive/2026-02-19/screenshots/manifest.json",
        help="Screenshot manifest json path",
    )
    parser.add_argument(
        "--output",
        default="reports/executive/2026-02-19/SyncHub_경영진_기획보고서_2026-02-19.pptx",
        help="Output pptx path",
    )
    args = parser.parse_args()

    context = _read_json(Path(args.context).resolve())
    manifest = _read_json(Path(args.manifest).resolve())
    output_path = Path(args.output).resolve()

    build_ppt(context, manifest, output_path)
    print(f"[ok] wrote pptx: {output_path}")


if __name__ == "__main__":
    main()
